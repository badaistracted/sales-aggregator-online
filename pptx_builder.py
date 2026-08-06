# pptx_builder.py
"""
Assembles the monthly management PowerPoint.
Uses python-pptx with a programmatic theme (no external template file needed).
"""
import io
import math
import calendar as cal
from datetime import datetime, timedelta
from collections import defaultdict

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Import chart builders from the same package
from chart_builder import chart_all_tenants_line, chart_single_tenant_line

# ── Slide dimensions (16:9 widescreen) ────────────────────────
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ── Brand colours ──────────────────────────────────────────────
C_NAVY      = RGBColor(0x1A, 0x3A, 0x5C)
C_BLUE      = RGBColor(0x28, 0x74, 0xA6)
C_TEAL      = RGBColor(0x1A, 0xBC, 0x9C)
C_AMBER     = RGBColor(0xF3, 0x9C, 0x12)
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT     = RGBColor(0xEC, 0xF0, 0xF1)
C_TEXT      = RGBColor(0x2C, 0x3E, 0x50)
C_MUTED     = RGBColor(0x7F, 0x8C, 0x8D)
C_RED       = RGBColor(0xE7, 0x4C, 0x3C)
C_GREEN     = RGBColor(0x27, 0xAE, 0x60)
C_PURPLE    = RGBColor(0x68, 0x5B, 0xC7)   # MoM % sticker colour


def _rgb_to_hex(rgb: RGBColor) -> str:
    return f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"


# ── Low-level helpers ──────────────────────────────────────────

def _add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    return shape


def _add_textbox(
    slide, text, left, top, width, height,
    font_size=12, bold=False, color=C_TEXT,
    align=PP_ALIGN.LEFT, wrap=True, italic=False,
):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size    = Pt(font_size)
    run.font.bold    = bold
    run.font.italic  = italic
    run.font.color.rgb = color
    run.font.name    = "Calibri"
    return txb


def _add_image(slide, img_buf, left, top, width, height):
    if img_buf is None:
        return
    img_buf.seek(0)
    slide.shapes.add_picture(img_buf, left, top, width, height)


def _add_bullet_textbox(
    slide, items: list, left, top, width, height,
    font_size=11, color=C_TEXT, bullet_char="▸",
):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"{bullet_char}  {item}"
        run.font.size      = Pt(font_size)
        run.font.color.rgb = color
        run.font.name      = "Calibri"
        p.space_before     = Pt(4)


def _header_bar(slide, title: str, subtitle: str = ""):
    """Dark navy header bar at top of every content slide."""
    _add_rect(slide, 0, 0, SLIDE_W, Inches(1.1), C_NAVY)
    _add_textbox(
        slide, title,
        left=Inches(0.4), top=Inches(0.1),
        width=Inches(10), height=Inches(0.65),
        font_size=22, bold=True, color=C_WHITE,
    )
    if subtitle:
        _add_textbox(
            slide, subtitle,
            left=Inches(0.4), top=Inches(0.72),
            width=Inches(9), height=Inches(0.35),
            font_size=11, color=RGBColor(0xA9, 0xCC, 0xE3),
        )
    _add_rect(slide, 0, Inches(1.1), SLIDE_W, Inches(0.045), C_TEAL)


def _footer_bar(slide, text: str = ""):
    """Light grey footer strip at the bottom of every slide."""
    _add_rect(slide, 0, SLIDE_H - Inches(0.32), SLIDE_W, Inches(0.32), C_LIGHT)
    if text:
        _add_textbox(
            slide, text,
            left=Inches(0.35), top=SLIDE_H - Inches(0.29),
            width=Inches(12), height=Inches(0.25),
            font_size=7.5, color=C_MUTED,
        )


def _kpi_box(slide, label, value, left, top,
             width=Inches(2.5), height=Inches(1.2),
             value_color=C_NAVY, bg_color=C_LIGHT):
    _add_rect(slide, left, top, width, height, bg_color)
    _add_textbox(
        slide, str(value),
        left=left + Inches(0.1), top=top + Inches(0.1),
        width=width - Inches(0.2), height=height * 0.55,
        font_size=18, bold=True, color=value_color,
        align=PP_ALIGN.CENTER,
    )
    _add_textbox(
        slide, label,
        left=left + Inches(0.1), top=top + height * 0.58,
        width=width - Inches(0.2), height=height * 0.42,
        font_size=9, color=C_MUTED, align=PP_ALIGN.CENTER,
    )


def _blank_slide(prs):
    """Return a new blank slide with white background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, C_WHITE)
    return slide


# ── Existing slide builders ────────────────────────────────────

def _slide_cover(prs, month_label: str, generated_at: str):
    layout = prs.slide_layouts[6]
    slide  = prs.slides.add_slide(layout)
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, C_NAVY)
    _add_rect(slide, 0, 0, Inches(0.18), SLIDE_H, C_TEAL)
    _add_textbox(
        slide, "MALL MANAGEMENT",
        left=Inches(0.5), top=Inches(1.5),
        width=Inches(10), height=Inches(0.7),
        font_size=14, color=RGBColor(0xA9, 0xCC, 0xE3), bold=True,
    )
    _add_textbox(
        slide, "K Square Monthly Performance Report",
        left=Inches(0.5), top=Inches(2.2),
        width=Inches(11), height=Inches(1.2),
        font_size=40, bold=True, color=C_WHITE,
    )
    _add_textbox(
        slide, month_label,
        left=Inches(0.5), top=Inches(3.4),
        width=Inches(8), height=Inches(0.8),
        font_size=28, color=C_TEAL, bold=True,
    )
    _add_rect(slide, Inches(0.5), Inches(4.25), Inches(5), Inches(0.04), C_TEAL)
    _add_textbox(
        slide, f"Generated: {generated_at}",
        left=Inches(0.5), top=Inches(4.4),
        width=Inches(8), height=Inches(0.4),
        font_size=10, color=C_MUTED,
    )
    _add_textbox(
        slide, "CONFIDENTIAL — Management Use Only",
        left=Inches(0.5), top=Inches(6.8),
        width=Inches(10), height=Inches(0.4),
        font_size=9, color=C_MUTED, italic=True,
    )


def _slide_executive_summary(prs, kpis: dict, llm_text: dict, month_label: str):
    layout = prs.slide_layouts[6]
    slide  = prs.slides.add_slide(layout)
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, C_WHITE)
    _header_bar(slide, "Executive Summary", month_label)

    def _fmt_idr(v):
        if v >= 1_000_000_000: return f"IDR {v/1_000_000_000:.2f}B"
        if v >= 1_000_000:     return f"IDR {v/1_000_000:.1f}M"
        return f"IDR {v:,.0f}"

    kpi_top = Inches(1.35)
    kpi_w   = Inches(2.4)
    kpi_h   = Inches(1.3)
    gap     = Inches(0.22)

    growth = kpis.get("sales_growth_pct")
    growth_label = (
        f"Total Sales  ({growth:+.1f}% MoM)" if growth is not None else "Total Sales"
    )
    v_color = (C_GREEN if growth and growth >= 0 else C_RED) if growth is not None else C_NAVY
    _kpi_box(slide, growth_label, _fmt_idr(kpis["total_sales"]),
             left=Inches(0.4), top=kpi_top, width=kpi_w, height=kpi_h, value_color=v_color)

    traffic = kpis["total_traffic"]
    _kpi_box(slide, "Total Visitors",
             f"{traffic:,.0f}" if traffic > 0 else "N/A",
             left=Inches(0.4) + kpi_w + gap, top=kpi_top, width=kpi_w, height=kpi_h)

    spv = kpis.get("sales_per_visitor")
    _kpi_box(slide, "Sales / Visitor",
             _fmt_idr(spv) if spv else "N/A",
             left=Inches(0.4) + 2 * (kpi_w + gap), top=kpi_top, width=kpi_w, height=kpi_h)

    _kpi_box(slide, "Active Tenants", str(kpis["total_tenants"]),
             left=Inches(0.4) + 3 * (kpi_w + gap), top=kpi_top, width=kpi_w, height=kpi_h)

    _kpi_box(slide, f"Top: {kpis['top_tenant_name']}", _fmt_idr(kpis["top_tenant_sales"]),
             left=Inches(0.4) + 4 * (kpi_w + gap), top=kpi_top,
             width=kpi_w, height=kpi_h, value_color=C_TEAL)

    _add_textbox(
        slide, llm_text.get("executive_summary", ""),
        left=Inches(0.4), top=Inches(2.85),
        width=Inches(12.4), height=Inches(2.6),
        font_size=12, color=C_TEXT, wrap=True,
    )

    _add_rect(slide, 0, SLIDE_H - Inches(0.35), SLIDE_W, Inches(0.35), C_LIGHT)
    _add_textbox(
        slide, f"Data as of {month_label}  |  Tenant count: {kpis['total_tenants']}",
        left=Inches(0.3), top=SLIDE_H - Inches(0.33),
        width=Inches(10), height=Inches(0.3),
        font_size=8, color=C_MUTED,
    )


def _slide_with_chart(
    prs, title: str, subtitle: str,
    chart_buf: io.BytesIO, notes_text: str,
    chart_left=Inches(0.3), chart_top=Inches(1.25),
    chart_w=Inches(8.5), chart_h=Inches(5.0),
):
    layout = prs.slide_layouts[6]
    slide  = prs.slides.add_slide(layout)
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, C_WHITE)
    _header_bar(slide, title, subtitle)
    _add_image(slide, chart_buf, left=chart_left, top=chart_top,
               width=chart_w, height=chart_h)

    panel_left = chart_left + chart_w + Inches(0.2)
    panel_w    = SLIDE_W - panel_left - Inches(0.2)
    if panel_w > Inches(1.0):
        _add_rect(slide, panel_left, Inches(1.25),
                  panel_w, Inches(5.5), RGBColor(0xF4, 0xF6, 0xF9))
        _add_textbox(
            slide, "Key Findings",
            left=panel_left + Inches(0.1), top=Inches(1.35),
            width=panel_w - Inches(0.2), height=Inches(0.4),
            font_size=10, bold=True, color=C_NAVY,
        )
        _add_textbox(
            slide, notes_text,
            left=panel_left + Inches(0.1), top=Inches(1.85),
            width=panel_w - Inches(0.2), height=Inches(4.8),
            font_size=10, color=C_TEXT, wrap=True,
        )

    _add_rect(slide, 0, SLIDE_H - Inches(0.3), SLIDE_W, Inches(0.3), C_LIGHT)
    return slide


def _slide_events(prs, events_data: dict, target_key: str, month_label: str):
    all_events = events_data.get("events_flat", [])
    if target_key:
        target_events = [
            e for e in all_events
            if str(e.get("date", "")).startswith(target_key)
        ]
    else:
        target_events = all_events

    target_events = sorted(
        target_events,
        key=lambda e: (
            str(e.get("date", "")),
            str(e.get("location", "")),
            str(e.get("event_name", "")),
        ),
    )

    def _new_event_slide(page_num=1, total_pages=1):
        layout = prs.slide_layouts[6]
        slide  = prs.slides.add_slide(layout)
        _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, C_WHITE)
        title = "Event Calendar"
        if total_pages > 1:
            title = f"Event Calendar ({page_num}/{total_pages})"
        _header_bar(slide, title, month_label)
        _add_textbox(
            slide, f"{len(target_events)} event(s) in this period",
            left=Inches(0.4), top=Inches(1.25),
            width=Inches(12), height=Inches(0.3),
            font_size=9, color=C_MUTED, italic=True,
        )
        return slide

    if not target_events:
        slide = _new_event_slide()
        _add_textbox(slide, "No events found for this period.",
                     left=Inches(0.4), top=Inches(1.8),
                     width=Inches(10), height=Inches(0.5),
                     font_size=14, color=C_MUTED)
        return

    rows_per_slide = 22
    chunks = [
        target_events[i:i + rows_per_slide]
        for i in range(0, len(target_events), rows_per_slide)
    ]
    total_pages = len(chunks)

    for page_idx, chunk in enumerate(chunks, start=1):
        slide = _new_event_slide(page_idx, total_pages)

        left   = Inches(0.35)
        top    = Inches(1.65)
        width  = Inches(12.65)
        height = Inches(5.55)
        rows   = len(chunk) + 1
        cols   = 3

        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
        table = table_shape.table
        table.columns[0].width = Inches(1.25)
        table.columns[1].width = Inches(1.75)
        table.columns[2].width = Inches(9.65)
        table.rows[0].height   = Inches(0.28)
        for r in range(1, rows):
            table.rows[r].height = Inches(0.24)

        def _set_cell(cell, text, font_size=7, bold=False,
                      color=C_TEXT, fill=None, align=PP_ALIGN.LEFT):
            if fill:
                cell.fill.solid()
                cell.fill.fore_color.rgb = fill
            tf = cell.text_frame
            tf.clear()
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = align
            run = p.add_run()
            run.text = str(text)
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = color
            run.font.name = "Calibri"
            cell.margin_left   = Inches(0.04)
            cell.margin_right  = Inches(0.04)
            cell.margin_top    = Inches(0.02)
            cell.margin_bottom = Inches(0.02)

        for c, h in enumerate(["Date", "Location", "Event"]):
            _set_cell(table.cell(0, c), h, font_size=7, bold=True,
                      color=C_WHITE, fill=C_NAVY, align=PP_ALIGN.CENTER)

        for i, evt in enumerate(chunk, start=1):
            raw_date = str(evt.get("date", ""))
            try:
                dt = datetime.strptime(raw_date, "%Y-%m-%d")
                date_label = dt.strftime("%d-%b")
                is_weekend = dt.weekday() >= 5
            except ValueError:
                date_label = raw_date
                is_weekend = False

            row_fill = (RGBColor(0xE8, 0xF5, 0xE9) if is_weekend
                        else (RGBColor(0xFF, 0xFF, 0xFF) if i % 2
                              else RGBColor(0xF7, 0xF9, 0xFC)))

            _set_cell(table.cell(i, 0), date_label, font_size=6.5,
                      bold=is_weekend, color=C_TEXT,
                      fill=row_fill, align=PP_ALIGN.CENTER)
            _set_cell(table.cell(i, 1), evt.get("location", "") or "-",
                      font_size=6.3, color=C_TEXT,
                      fill=row_fill, align=PP_ALIGN.CENTER)
            _set_cell(table.cell(i, 2), evt.get("event_name", "") or "-",
                      font_size=6.3, color=C_TEXT,
                      fill=row_fill, align=PP_ALIGN.LEFT)

        _footer_bar(
            slide,
            f"Events shown: {(page_idx - 1) * rows_per_slide + 1}"
            f"–{(page_idx - 1) * rows_per_slide + len(chunk)}"
            f" of {len(target_events)}",
        )


def _slide_event_impact(prs, kpis: dict, target_key: str, month_label: str):
    event_impact    = kpis.get("event_impact", [])
    avg_event       = kpis.get("avg_event_day_sales", 0)
    avg_non_event   = kpis.get("avg_non_event_sales", 0)
    event_days      = kpis.get("event_days_count", 0)
    non_event_days  = kpis.get("non_event_days_count", 0)

    if not event_impact:
        return

    layout = prs.slide_layouts[6]
    slide  = prs.slides.add_slide(layout)
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, C_WHITE)
    _header_bar(slide, "Event Impact Analysis", month_label)

    def _fmt(v):
        if v >= 1_000_000_000: return f"IDR {v/1_000_000_000:.1f}B"
        if v >= 1_000_000:     return f"IDR {v/1_000_000:.1f}M"
        if v >= 1_000:         return f"IDR {v/1_000:.0f}K"
        return f"IDR {v:,.0f}"

    kpi_top = Inches(1.25)
    kpi_w   = Inches(3.0)
    kpi_h   = Inches(1.0)
    gap     = Inches(0.25)

    _kpi_box(slide, f"Avg Sales — Event Days ({event_days}d)", _fmt(avg_event),
             left=Inches(0.35), top=kpi_top, width=kpi_w, height=kpi_h,
             value_color=C_TEAL, bg_color=C_LIGHT)

    _kpi_box(slide, f"Avg Sales — No Events ({non_event_days}d)", _fmt(avg_non_event),
             left=Inches(0.35) + kpi_w + gap, top=kpi_top,
             width=kpi_w, height=kpi_h, value_color=C_NAVY, bg_color=C_LIGHT)

    if avg_non_event > 0:
        uplift_pct   = (avg_event - avg_non_event) / avg_non_event * 100
        uplift_label = f"{uplift_pct:+.1f}%"
        uplift_color = C_GREEN if uplift_pct > 0 else C_RED
    else:
        uplift_label = "N/A"
        uplift_color = C_MUTED

    _kpi_box(slide, "Event Day Uplift", uplift_label,
             left=Inches(0.35) + 2 * (kpi_w + gap), top=kpi_top,
             width=kpi_w, height=kpi_h, value_color=uplift_color, bg_color=C_LIGHT)

    best = event_impact[0] if event_impact else None
    if best:
        try:
            dt = datetime.strptime(best["date"], "%Y-%m-%d")
            best_label = dt.strftime("%d %b")
        except ValueError:
            best_label = best["date"]
        _kpi_box(slide, f"Best Day: {best_label}", _fmt(best["sales"]),
                 left=Inches(0.35) + 3 * (kpi_w + gap), top=kpi_top,
                 width=kpi_w, height=kpi_h, value_color=C_AMBER, bg_color=C_LIGHT)

    top_events  = event_impact[:15]
    table_shape = slide.shapes.add_table(
        len(top_events) + 1, 5,
        Inches(0.35), Inches(2.5), Inches(12.65), Inches(4.5),
    )
    table = table_shape.table
    table.columns[0].width = Inches(0.45)
    table.columns[1].width = Inches(1.15)
    table.columns[2].width = Inches(1.65)
    table.columns[3].width = Inches(1.40)
    table.columns[4].width = Inches(8.00)

    def _set_cell(cell, text, font_size=7, bold=False,
                  color=C_TEXT, fill=None, align=PP_ALIGN.LEFT):
        if fill:
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill
        tf = cell.text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = str(text)
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Calibri"
        cell.margin_left = cell.margin_right = Inches(0.04)
        cell.margin_top  = cell.margin_bottom = Inches(0.02)

    for c, h in enumerate(["#", "Date", "Sales (IDR)", "vs Avg", "Events"]):
        _set_cell(table.cell(0, c), h, font_size=7, bold=True,
                  color=C_WHITE, fill=C_NAVY, align=PP_ALIGN.CENTER)

    table.rows[0].height = Inches(0.28)
    for i, ei in enumerate(top_events, start=1):
        try:
            dt = datetime.strptime(ei["date"], "%Y-%m-%d")
            date_label = dt.strftime("%a, %d %b")
            is_weekend = dt.weekday() >= 5
        except ValueError:
            date_label = ei["date"]
            is_weekend = False

        row_fill = (RGBColor(0xE8, 0xF5, 0xE9) if is_weekend
                    else (RGBColor(0xFF, 0xFF, 0xFF) if i % 2
                          else RGBColor(0xF7, 0xF9, 0xFC)))

        uplift_text  = f"{ei['uplift_pct']:+.0f}%"
        uplift_color = C_GREEN if ei["uplift_pct"] > 0 else (C_RED if ei["uplift_pct"] < 0 else C_MUTED)
        events_text  = " · ".join(ei.get("events", [])[:4])
        if len(ei.get("events", [])) > 4:
            events_text += f" +{len(ei['events']) - 4} more"

        _set_cell(table.cell(i, 0), str(i), font_size=7, bold=(i <= 3),
                  color=C_AMBER if i <= 3 else C_TEXT,
                  fill=row_fill, align=PP_ALIGN.CENTER)
        _set_cell(table.cell(i, 1), date_label, font_size=6.5,
                  bold=is_weekend, color=C_TEXT, fill=row_fill, align=PP_ALIGN.CENTER)
        _set_cell(table.cell(i, 2), _fmt(ei["sales"]), font_size=7,
                  bold=True, color=C_NAVY, fill=row_fill, align=PP_ALIGN.RIGHT)
        _set_cell(table.cell(i, 3), uplift_text, font_size=7,
                  bold=True, color=uplift_color, fill=row_fill, align=PP_ALIGN.CENTER)
        _set_cell(table.cell(i, 4), events_text, font_size=6.3,
                  color=C_TEXT, fill=row_fill, align=PP_ALIGN.LEFT)
        table.rows[i].height = Inches(0.27)

    _footer_bar(
        slide,
        f"Top {len(top_events)} event day(s) by total sales. "
        f"Uplift vs average non-event day ({_fmt(avg_non_event)}).",
    )


def _slide_recommendations(prs, recommendations: list, month_label: str):
    layout = prs.slide_layouts[6]
    slide  = prs.slides.add_slide(layout)
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, C_NAVY)
    _add_textbox(
        slide, "Recommendations & Next Steps",
        left=Inches(0.5), top=Inches(0.4),
        width=Inches(12), height=Inches(0.7),
        font_size=24, bold=True, color=C_WHITE,
    )
    _add_textbox(slide, month_label,
                 left=Inches(0.5), top=Inches(1.1),
                 width=Inches(8), height=Inches(0.4),
                 font_size=12, color=C_TEAL)
    _add_rect(slide, Inches(0.5), Inches(1.55), Inches(12.3), Inches(0.04), C_TEAL)

    card_colors = [
        RGBColor(0x1B, 0x4F, 0x72), RGBColor(0x1A, 0x53, 0x76),
        RGBColor(0x1A, 0x5C, 0x88), RGBColor(0x1F, 0x61, 0x8D),
    ]
    card_h = Inches(1.15)
    card_gap = Inches(0.18)
    card_top_start = Inches(1.75)

    for i, rec in enumerate(recommendations[:4]):
        top   = card_top_start + i * (card_h + card_gap)
        color = card_colors[i % len(card_colors)]
        _add_rect(slide, Inches(0.5), top, Inches(12.3), card_h, color)
        _add_textbox(
            slide, f"0{i+1}",
            left=Inches(0.6), top=top + Inches(0.25),
            width=Inches(0.6), height=Inches(0.65),
            font_size=22, bold=True, color=C_TEAL,
        )
        _add_textbox(
            slide, rec,
            left=Inches(1.3), top=top + Inches(0.12),
            width=Inches(11.3), height=card_h - Inches(0.24),
            font_size=13, color=C_WHITE, wrap=True,
        )

    _add_textbox(
        slide, "This report was generated automatically. Data subject to final audit.",
        left=Inches(0.5), top=Inches(7.1),
        width=Inches(12), height=Inches(0.3),
        font_size=8, color=C_MUTED, italic=True,
    )


# ══════════════════════════════════════════════════════════════
#  NEW: Percentage-Summary slide builders
# ══════════════════════════════════════════════════════════════

def _slide_all_tenants_trend(prs, pct_data_parsed: dict, month_label: str):
    """
    Slide: Full-width all-tenants line chart.
    Uses the white-background chart_all_tenants_line() function.
    """
    tenants_data = pct_data_parsed.get("tenants", {})
    all_months   = pct_data_parsed.get("all_months", [])
    pct_data     = pct_data_parsed.get("percentage", {})
    target_key   = pct_data_parsed.get("pct_months", {}).get("to", "")

    if not tenants_data or not all_months:
        return

    chart_buf = chart_all_tenants_line(tenants_data, all_months, pct_data, target_key)
    if chart_buf is None:
        return

    slide = _blank_slide(prs)
    _header_bar(slide, "Monthly Sales Trend — All Tenants", month_label)

    # Period label (top-right, compact)
    if all_months:
        period_str = f"{all_months[0]} → {all_months[-1]}"
        _add_textbox(
            slide, period_str,
            left=Inches(9.5), top=Inches(0.72),
            width=Inches(3.6), height=Inches(0.32),
            font_size=9, color=RGBColor(0xA9, 0xCC, 0xE3),
            align=PP_ALIGN.RIGHT,
        )

    # Full-width chart (tall, leaves room for legend below)
    _add_image(
        slide, chart_buf,
        left=Inches(0.2), top=Inches(1.25),
        width=Inches(12.9), height=Inches(5.85),
    )

    _footer_bar(
        slide,
        f"{len(tenants_data)} tenant(s) · "
        f"Circle = MoM % vs previous month · "
        f"Target: {month_label}",
    )


def _slide_tenant_grid(prs, pct_data_parsed: dict, month_label: str):
    """
    Slides: 6 tenant mini-charts per slide in a 3-column × 2-row grid.
    Tenants sorted by latest-month sales descending.
    All tenants included (Q21).
    """
    tenants_data = pct_data_parsed.get("tenants", {})
    all_months   = pct_data_parsed.get("all_months", [])
    pct_data     = pct_data_parsed.get("percentage", {})
    target_key   = pct_data_parsed.get("pct_months", {}).get("to", "")

    if not tenants_data or not all_months:
        return

    # Sort tenants by latest month sales descending
    def _sort_key(t):
        m = tenants_data[t].get("monthly", {})
        if target_key and target_key in m:
            return m[target_key]
        vals = [v for v in m.values() if v is not None]
        return vals[-1] if vals else 0

    tenant_names = sorted(tenants_data.keys(), key=_sort_key, reverse=True)

    # Batch into groups of 6
    batch_size = 6
    batches    = [
        tenant_names[i:i + batch_size]
        for i in range(0, len(tenant_names), batch_size)
    ]
    total_pages = len(batches)

    # Grid layout constants — 3 cols × 2 rows
    # Slide content area: x 0.2–13.13, y 1.25–7.18 (below header, above footer)
    CHART_W  = Inches(4.15)
    CHART_H  = Inches(2.75)
    H_GAP    = Inches(0.18)
    V_GAP    = Inches(0.18)
    GRID_LEFT = Inches(0.2)
    GRID_TOP  = Inches(1.3)

    for batch_idx, batch in enumerate(batches):
        slide = _blank_slide(prs)
        _header_bar(
            slide,
            f"Tenant Performance ({batch_idx + 1}/{total_pages})",
            month_label,
        )

        for i, tenant in enumerate(batch):
            col = i % 3
            row = i // 3

            left = GRID_LEFT + col * (CHART_W + H_GAP)
            top  = GRID_TOP  + row * (CHART_H + V_GAP)

            monthly = tenants_data[tenant].get("monthly", {})
            pct_val = pct_data.get(tenant, {}).get("pct") if pct_data else None

            img_buf = chart_single_tenant_line(
                tenant_name=tenant,
                monthly_data=monthly,
                all_months=all_months,
                pct_val=pct_val,
                target_key=target_key,
            )
            if img_buf:
                _add_image(slide, img_buf, left, top, CHART_W, CHART_H)

        _footer_bar(
            slide,
            f"Tenants {batch_idx * batch_size + 1}"
            f"–{min((batch_idx + 1) * batch_size, len(tenant_names))}"
            f" of {len(tenant_names)} · "
            f"Sorted by {target_key or 'latest'} sales",
        )


def _slide_summary_text(prs, pct_data_parsed: dict, month_label: str):
    """
    Slide: Raw Indonesian summary paragraph text extracted from the
    PERSENTASE % sheet.  Text is displayed as-is (Q15).
    """
    summary_text = pct_data_parsed.get("summary_text", [])
    if not summary_text:
        return

    slide = _blank_slide(prs)
    _header_bar(slide, "Analisis Penjualan", month_label)

    # Accent label
    pct_months = pct_data_parsed.get("pct_months", {})
    if pct_months.get("from") and pct_months.get("to"):
        period_str = f"Perbandingan: {pct_months['from']} vs {pct_months['to']}"
        _add_textbox(
            slide, period_str,
            left=Inches(0.45), top=Inches(1.2),
            width=Inches(9), height=Inches(0.32),
            font_size=9, color=C_MUTED, italic=True,
        )

    # Build a single multi-paragraph textbox
    # Content area: x 0.45, y 1.6, width 12.4, height 5.5
    txb = slide.shapes.add_textbox(
        Inches(0.45), Inches(1.6), Inches(12.4), Inches(5.5)
    )
    tf = txb.text_frame
    tf.word_wrap = True

    for para_idx, paragraph in enumerate(summary_text):
        paragraph = paragraph.strip()
        if not paragraph:
            continue

        p = tf.paragraphs[0] if para_idx == 0 else tf.add_paragraph()
        p.space_before = Pt(0 if para_idx == 0 else 10)

        run = p.add_run()
        run.text           = paragraph
        run.font.name      = "Calibri"
        run.font.color.rgb = C_TEXT

        # Numbered bullet points get slightly smaller font + indent
        if paragraph and paragraph[0].isdigit() and len(paragraph) > 1 and paragraph[1] in ".):":
            run.font.size = Pt(10.5)
            p.level       = 1
        else:
            run.font.size = Pt(11.5)
            run.font.bold = True   # introductory sentence stands out

    _footer_bar(slide, "Teks ringkasan diambil langsung dari laporan PERSENTASE %.")


# ── Main assembler ─────────────────────────────────────────────

def build_pptx(
kpis,
llm_text,
charts,
month_label,
target_key=None,
events_data=None,
) -> io.BytesIO:
"""
Assemble the full presentation.

charts keys:
    monthly_sales   — bar chart
    top_tenants     — horizontal bar
    traffic         — dual-axis line/bar
    daily_sales     — daily line chart
"""
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# If server runs in UTC and you want local time = UTC+5
generated_at = (datetime.utcnow() + timedelta(hours=7)).strftime("%d %B %Y, %H:%M")

# ── Slide 1: Cover ────────────────────────────────────────
_slide_cover(prs, month_label, generated_at)

# ── Slide 2: Executive Summary ────────────────────────────
_slide_executive_summary(prs, kpis, llm_text, month_label)

# ── Slide 3: Total Monthly Sales ─────────────────────────
if charts.get("monthly_sales"):
    _slide_with_chart(
        prs,
        title    = "Monthly Sales Performance",
        subtitle = month_label,
        chart_buf   = charts["monthly_sales"],
        notes_text  = llm_text.get("sales_slide_notes", ""),
        chart_left  = Inches(0.3),
        chart_top   = Inches(1.25),
        chart_w     = Inches(8.8),
        chart_h     = Inches(5.5),
    )

# ── Slide 4: Top Tenants ──────────────────────────────────
if charts.get("top_tenants"):
    _slide_with_chart(
        prs,
        title    = "Top Tenant Performance",
        subtitle = month_label,
        chart_buf   = charts["top_tenants"],
        notes_text  = llm_text.get("tenant_slide_notes", ""),
        chart_left  = Inches(0.3),
        chart_top   = Inches(1.25),
        chart_w     = Inches(8.8),
        chart_h     = Inches(5.5),
    )

# ── Slide 5: Traffic ──────────────────────────────────────
if charts.get("traffic"):
    _slide_with_chart(
        prs,
        title    = "Visitor Traffic & Spend Efficiency",
        subtitle = month_label,
        chart_buf   = charts["traffic"],
        notes_text  = llm_text.get("traffic_slide_notes", ""),
        chart_left  = Inches(0.3),
        chart_top   = Inches(1.25),
        chart_w     = Inches(8.8),
        chart_h     = Inches(5.5),
    )

# ── Slide 6: Daily Sales ──────────────────────────────────
if charts.get("daily_sales"):
    _slide_with_chart(
        prs,
        title    = "Daily Sales Pattern",
        subtitle = month_label,
        chart_buf   = charts["daily_sales"],
        notes_text  = llm_text.get("daily_slide_notes", ""),
        chart_left  = Inches(0.3),
        chart_top   = Inches(1.25),
        chart_w     = Inches(8.8),
        chart_h     = Inches(5.5),
    )

# ── Slide 7: Events Calendar ──────────────────────────────
if events_data and events_data.get("events_flat"):
    _slide_events(prs, events_data, target_key, month_label)

# ── Slide 8: Event Impact Analysis ────────────────────────
if kpis.get("event_impact"):
    _slide_event_impact(prs, kpis, target_key, month_label)

# ── Slide 9: Recommendations ─────────────────────────────
_slide_recommendations(
    prs,
    recommendations = llm_text.get("recommendations", []),
    month_label     = month_label,
)

# ── Save to BytesIO ───────────────────────────────────────
buf = io.BytesIO()
prs.save(buf)
buf.seek(0)
return buf

